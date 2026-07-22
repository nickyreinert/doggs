"""Document parsing utilities (backend/processor.py).

Copied from `app/processor.py` and adjusted package layout.
"""
from __future__ import annotations
import os
from typing import Optional
from typing import Tuple

try:
    from unstructured.partition.auto import partition
except Exception:  # pragma: no cover - optional dependency
    partition = None


def extract_text(path: str) -> str:
    if partition is not None:
        elements = partition(filename=path)
        return "\n\n".join(str(e) for e in elements)

    ext = os.path.splitext(path)[1].lower()
    if ext in {".txt", ".md"}:
        return _read_text(path)
    if ext == ".pdf":
        # Detect whether the PDF has an embedded text layer. Prefer
        # extracting that text (fast + accurate). Only run OCR when the
        # PDF appears to be image/scanned only or when a text extractor
        # is not available.
        has_text = _pdf_has_text_layer(path)
        if has_text is True:
            txt = _extract_pdf_text_with_pypdf(path)
            if txt and txt.strip():
                return txt
            # Unexpected: fall back to OCR
        if has_text is False:
            return _extract_pdf_with_ocr(path)
        # has_text is None (pypdf unavailable or error) — safe fallback to OCR
        return _extract_pdf_with_ocr(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext in {".pptx"}:
        return _extract_pptx(path)

    raise RuntimeError("Install 'unstructured' or add a handler for this file type")


def _read_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _extract_docx(path: str) -> str:
    try:
        from docx import Document
    except Exception as e:  # pragma: no cover - optional dependency
        raise RuntimeError("python-docx required for .docx parsing") from e

    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text]
    return "\n\n".join(parts)


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except Exception as e:  # pragma: no cover - optional dependency
        raise RuntimeError("python-pptx required for .pptx parsing") from e

    pres = Presentation(path)
    texts: list[str] = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    texts.append(t)
    return "\n\n".join(texts)


def _pdf_has_text_layer(path: str, char_threshold_per_page: int = 20, page_ratio: float = 0.3) -> Optional[bool]:
    """Return True if PDF likely has a text layer, False if likely image-only.

    Returns None when a PDF parser is not available or an error occurred.
    The heuristic counts pages with >= `char_threshold_per_page` characters
    and declares a text PDF when the fraction of such pages >= `page_ratio`.
    """
    try:
        from pypdf import PdfReader
    except Exception:
        return None

    try:
        reader = PdfReader(path)
        pages = reader.pages
        if not pages:
            return False
        pages_with_text = 0
        for p in pages:
            try:
                t = p.extract_text() or ""
            except Exception:
                t = ""
            if len(t.strip()) >= char_threshold_per_page:
                pages_with_text += 1
        return (pages_with_text / len(pages)) >= page_ratio
    except Exception:
        return None


def _extract_pdf_text_with_pypdf(path: str) -> str:
    try:
        from pypdf import PdfReader
    except Exception as e:
        raise RuntimeError("pypdf required for PDF text extraction") from e

    try:
        reader = PdfReader(path)
        parts: list[str] = []
        for p in reader.pages:
            try:
                t = p.extract_text() or ""
            except Exception:
                t = ""
            if t.strip():
                parts.append(t)
        return "\n\n".join(parts)
    except Exception as e:
        raise RuntimeError("Failed reading PDF with pypdf") from e


def _extract_pdf_with_ocr(path: str) -> str:
    try:
        from pdf2image import convert_from_path
        import pytesseract
    except Exception as e:  # pragma: no cover - optional dependency
        raise RuntimeError("pdf2image + pytesseract required for PDF OCR") from e

    images = convert_from_path(path)
    texts: list[str] = []
    for img in images:
        texts.append(pytesseract.image_to_string(img))
    return "\n\n".join(t for t in texts if t.strip())


def _summarize_text(text: str, max_chars: int = 800) -> str:
    """Very small heuristic summary: return the first paragraph up to max_chars."""
    if not text:
        return ""
    parts = [p.strip() for p in text.split("\n\n") if p.strip()]
    if not parts:
        return text[:max_chars]
    summary = parts[0]
    if len(summary) > max_chars:
        return summary[:max_chars].rsplit(" ", 1)[0] + "..."
    return summary


def _get_metadata(path: str) -> dict:
    st = os.stat(path)
    meta: dict = {
        "path": path,
        "size": st.st_size,
        "mtime": int(st.st_mtime),
        "ctime": int(st.st_ctime),
    }
    # try format-specific metadata
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".docx":
            from docx import Document

            doc = Document(path)
            props = doc.core_properties
            meta.update({k: getattr(props, k, None) for k in ("author", "title", "created", "last_modified_by")})
        elif ext == ".pptx":
            from pptx import Presentation

            pres = Presentation(path)
            props = pres.core_properties
            meta.update({k: getattr(props, k, None) for k in ("author", "title", "created", "last_modified_by")})
        elif ext == ".pdf":
            try:
                from pypdf import PdfReader

                r = PdfReader(path)
                if hasattr(r, "metadata") and r.metadata:
                    for k, v in r.metadata.items():
                        if v:
                            meta.setdefault("pdf_metadata", {})[k] = str(v)
            except Exception:
                pass
    except Exception:
        # non-fatal; return basic fs metadata
        pass

    return meta


def _cut_text(text: str, max_width: int) -> str:
    if not text or max_width <= 0:
        return text
    parts: list[str] = []
    for para in text.split("\n\n"):
        para = para.strip()
        if not para:
            continue
        if len(para) <= max_width:
            parts.append(para)
            continue
        # cut into chunks at word boundaries
        start = 0
        L = len(para)
        while start < L:
            end = min(start + max_width, L)
            if end < L:
                # try to backtrack to last space
                bs = para.rfind(" ", start, end)
                if bs > start:
                    end = bs
            parts.append(para[start:end].strip())
            start = end + 1
    return "\n\n".join(parts)


def _parse_structure_from_text(text: str) -> dict:
    """Very small heuristic parser: identify headings, lists and body paragraphs."""
    headers: list[dict] = []
    lists: list[list[str]] = []
    paragraphs: list[str] = []

    for para in [p.strip() for p in text.split("\n\n") if p.strip()]:
        # list detection
        lines = [l.strip() for l in para.splitlines() if l.strip()]
        if all(l.startswith(('-', '*', '•')) or l[:2].isdigit() for l in lines) and len(lines) > 1:
            # treat as list
            clean = [l.lstrip('-*• ').lstrip('0123456789. ') for l in lines]
            lists.append(clean)
            continue

        # explicit markdown-like headings
        if para.startswith('#'):
            level = len(para) - len(para.lstrip('#'))
            text_h = para.lstrip('#').strip()
            headers.append({"level": level, "text": text_h})
            continue

        # ALL CAPS short line as heading heuristic
        single_line = para.splitlines()[0]
        words = single_line.split()
        if len(words) <= 10 and single_line.upper() == single_line and len(single_line) < 200:
            headers.append({"level": 1, "text": single_line})
            # remaining lines if any become paragraph
            rest = "\n".join(lines[1:]).strip()
            if rest:
                paragraphs.append(rest)
            continue

        paragraphs.append(para)

    return {"headers": headers, "lists": lists, "paragraphs": paragraphs}


def extract_structured(path: str, max_width: int = 200) -> dict:
    """Return structured data for a document: metadata, headers, lists, paragraphs.

    Uses format-specific metadata when available. Attempts richer parsing for
    `.docx` and `.pptx` using their libraries; falls back to heuristics on raw text.
    """
    meta = _get_metadata(path)

    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".docx":
            try:
                from docx import Document

                doc = Document(path)
                headers: list[dict] = []
                lists: list[list[str]] = []
                paras: list[str] = []
                for p in doc.paragraphs:
                    text = p.text.strip()
                    if not text:
                        continue
                    sty = getattr(p, "style", None)
                    name = getattr(sty, "name", "") if sty is not None else ""
                    if name.lower().startswith("heading"):
                        # extract level from 'Heading 1'
                        try:
                            level = int(name.split()[-1])
                        except Exception:
                            level = 1
                        headers.append({"level": level, "text": text})
                        continue
                    if name.lower().startswith("list") or text.startswith(('-', '*', '•')):
                        lists.append([text.lstrip('-*• ').strip()])
                        continue
                    paras.append(text)
                body = {"headers": headers, "lists": lists, "paragraphs": paras}
                # cut text
                body = {k: ([_cut_text(x, max_width) for x in v] if isinstance(v, list) else _cut_text(v, max_width)) for k, v in body.items()}
                return {"meta": meta, "structure": body}
            except Exception:
                pass

        if ext == ".pptx":
            try:
                from pptx import Presentation

                pres = Presentation(path)
                headers = []
                lists = []
                paras = []
                for slide in pres.slides:
                    for shape in slide.shapes:
                        if not hasattr(shape, "text"):
                            continue
                        text = shape.text.strip()
                        if not text:
                            continue
                        # python-pptx gives paragraph.level on text_frame paragraphs
                        try:
                            for p in shape.text_frame.paragraphs:
                                lvl = getattr(p, "level", 0)
                                t = "".join(r.text for r in p.runs).strip()
                                if not t:
                                    continue
                                if lvl == 0 and len(t.split()) <= 8:
                                    headers.append({"level": 1, "text": t})
                                elif t.startswith(('-', '*', '•')):
                                    lists.append([t.lstrip('-*• ').strip()])
                                else:
                                    paras.append(t)
                        except Exception:
                            paras.append(text)
                body = {"headers": headers, "lists": lists, "paragraphs": paras}
                body = {k: ([_cut_text(x, max_width) for x in v] if isinstance(v, list) else _cut_text(v, max_width)) for k, v in body.items()}
                return {"meta": meta, "structure": body}
            except Exception:
                pass

        # Generic fallback: extract raw text then parse heuristically
        text = extract_text(path)
        text = _cut_text(text, max_width)
        parsed = _parse_structure_from_text(text)
        return {"meta": meta, "structure": parsed}
    except Exception as e:
        return {"meta": meta, "error": str(e)}


def main(argv: Optional[list[str]] = None) -> int:
    import argparse
    import json

    parser = argparse.ArgumentParser(
        prog="processor",
        description="Extract text from documents. Uses 'unstructured' if available, otherwise falls back to format-specific handlers and OCR for PDFs.",
    )
    parser.add_argument("paths", nargs="*", help="One or more file paths to parse. If omitted, `--config` or $CONFIG_PATH will be used to determine target_path to scan.")
    parser.add_argument("--config", type=str, help="Path to config.json. If provided and no paths, processor will scan `scan_settings.target_path` from the config.")
    parser.add_argument("--head", type=int, default=0, help="Print first N characters of extracted text")
    parser.add_argument("--lines", type=int, default=0, help="Print first N lines of extracted text")
    parser.add_argument("--save", type=str, help="Save extracted text to given directory (creates files with .txt)")
    parser.add_argument("--json", action="store_true", help="Print JSON objects with {path, text, summary}")
    parser.add_argument("--structured", action="store_true", help="Print structured JSON {meta, structure}")
    parser.add_argument("--max-width", type=int, default=200, help="Max chars per paragraph chunk when structuring")
    parser.add_argument("--summary-chars", type=int, default=800, help="Max chars for heuristic summary when using --json or --summary")

    args = parser.parse_args(argv)

    # Resolve config path from arg or env
    config_path = args.config or os.getenv("CONFIG_PATH")
    config: dict | None = None
    if not args.paths and config_path:
        try:
            import json as _json

            with open(config_path, "r", encoding="utf-8") as cf:
                config = _json.load(cf)
        except Exception as e:
            print(f"ERROR reading config {config_path}: {e}")
            return 2

        # if config provided and no explicit paths, derive paths from config.scan_settings.target_path
        scan_settings = config.get("scan_settings", {})
        target = scan_settings.get("target_path")
        allowed = set(ext.lower() for ext in scan_settings.get("allowed_extensions", []))
        if not target:
            print("No target_path in config scan_settings")
            return 2
        # Walk the target and collect matching files
        collected: list[str] = []
        for dirpath, _, filenames in os.walk(target):
            for fn in filenames:
                if not allowed or os.path.splitext(fn)[1].lower() in allowed:
                    collected.append(os.path.join(dirpath, fn))
        args.paths = collected

    out_dir = None
    if args.save:
        out_dir = args.save
        os.makedirs(out_dir, exist_ok=True)

    if not args.paths:
        print("No input paths provided; nothing to do.")
        return 0

    for p in args.paths:
        try:
            text = extract_text(p)
        except Exception as e:
            print(f"ERROR parsing {p}: {e}")
            continue

        if args.structured:
            structured = extract_structured(p, max_width=args.max_width)
            print(json.dumps(structured, ensure_ascii=False))
            continue

        if args.json:
            summary = _summarize_text(text, max_chars=args.summary_chars)
            obj = {"path": p, "text": text, "summary": summary}
            print(json.dumps(obj, ensure_ascii=False))
            continue

        if args.head > 0:
            print(f"--- {p} (first {args.head} chars) ---")
            print(text[: args.head])
        elif args.lines > 0:
            print(f"--- {p} (first {args.lines} lines) ---")
            lines = text.splitlines()
            for L in lines[: args.lines]:
                print(L)
        else:
            print(f"--- {p} (full) ---")
            print(text)

        if out_dir:
            bn = os.path.basename(p)
            fn = os.path.splitext(bn)[0] + ".txt"
            out_path = os.path.join(out_dir, fn)
            try:
                with open(out_path, "w", encoding="utf-8") as f:
                    f.write(text)
            except Exception as e:
                print(f"ERROR saving to {out_path}: {e}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
